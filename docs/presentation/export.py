#!/usr/bin/env python3
"""
export.py – Convert presentation.html to PDF and/or PowerPoint (.pptx).

Usage
-----
    python export.py                      # both PDF and PPTX
    python export.py --format pdf         # PDF only
    python export.py --format pptx        # PPTX only
    python export.py --out my-deck        # custom output stem (no extension)

Requirements
------------
    pip install playwright python-pptx pillow
    playwright install chromium

How it works
------------
1. Playwright opens presentation.html in a headless Chromium browser
   at 1280 × 720 (16:9).
2. JavaScript is used to activate each slide one at a time.
3. A full-viewport screenshot is taken of every slide.
4. PDF  – screenshots are stitched into a single PDF via Pillow.
5. PPTX – each screenshot becomes a full-bleed image slide via python-pptx.
"""

import argparse
import os
import sys
import tempfile
from pathlib import Path

# ---------------------------------------------------------------------------
# Argument parsing
# ---------------------------------------------------------------------------

parser = argparse.ArgumentParser(description="Export presentation.html to PDF/PPTX")
parser.add_argument(
    "--format",
    choices=["pdf", "pptx", "both"],
    default="both",
    help="Output format (default: both)",
)
parser.add_argument(
    "--out",
    default="presentation",
    help="Output file stem without extension (default: presentation)",
)
parser.add_argument(
    "--width",
    type=int,
    default=1280,
    help="Viewport width in pixels (default: 1280)",
)
parser.add_argument(
    "--height",
    type=int,
    default=720,
    help="Viewport height in pixels (default: 720)",
)
args = parser.parse_args()

# ---------------------------------------------------------------------------
# Resolve paths
# ---------------------------------------------------------------------------

script_dir = Path(__file__).parent.resolve()
html_file = script_dir / "presentation.html"

if not html_file.exists():
    print(f"ERROR: {html_file} not found", file=sys.stderr)
    sys.exit(1)

html_url = html_file.as_uri()

# ---------------------------------------------------------------------------
# Capture slides with Playwright
# ---------------------------------------------------------------------------

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    print(
        "ERROR: playwright is not installed.\n"
        "Run:  pip install playwright && playwright install chromium",
        file=sys.stderr,
    )
    sys.exit(1)

print(f"Opening {html_url} …")

slide_images: list[Path] = []

with sync_playwright() as pw:
    browser = pw.chromium.launch()
    page = browser.new_page(viewport={"width": args.width, "height": args.height})
    page.goto(html_url, wait_until="networkidle")

    # Count slides
    slide_count: int = page.evaluate("() => document.querySelectorAll('.slide').length")
    print(f"Found {slide_count} slides")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)

        for i in range(slide_count):
            # Activate slide i via the jumpTo function defined in the HTML
            page.evaluate(f"() => jumpTo({i})")
            page.wait_for_timeout(120)  # let CSS transitions settle

            img_path = tmp_path / f"slide_{i:03d}.png"
            page.screenshot(path=str(img_path), full_page=False)
            slide_images.append(img_path)
            print(f"  Captured slide {i + 1}/{slide_count}")

        # -----------------------------------------------------------------------
        # Export PDF
        # -----------------------------------------------------------------------
        if args.format in ("pdf", "both"):
            try:
                from PIL import Image
            except ImportError:
                print(
                    "ERROR: Pillow is not installed.\nRun:  pip install pillow",
                    file=sys.stderr,
                )
                sys.exit(1)

            pdf_path = script_dir / f"{args.out}.pdf"
            images = [Image.open(p).convert("RGB") for p in slide_images]
            first, rest = images[0], images[1:]
            first.save(pdf_path, save_all=True, append_images=rest)
            print(f"PDF saved → {pdf_path}")

        # -----------------------------------------------------------------------
        # Export PPTX
        # -----------------------------------------------------------------------
        if args.format in ("pptx", "both"):
            try:
                from pptx import Presentation
                from pptx.util import Emu
            except ImportError:
                print(
                    "ERROR: python-pptx is not installed.\nRun:  pip install python-pptx",
                    file=sys.stderr,
                )
                sys.exit(1)

            prs = Presentation()

            # Set slide size to 16:9 (1280×720 px → scaled to standard 10 × 5.625 inches)
            slide_width_emu = Emu(int(9144000 * args.width / 1280))   # 9 144 000 EMU = 10 in
            slide_height_emu = Emu(int(5143500 * args.height / 720))  # 5 143 500 EMU = 5.625 in
            prs.slide_width = slide_width_emu
            prs.slide_height = slide_height_emu

            blank_layout = prs.slide_layouts[6]  # completely blank layout

            for img_path in slide_images:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(img_path),
                    left=0,
                    top=0,
                    width=slide_width_emu,
                    height=slide_height_emu,
                )

            pptx_path = script_dir / f"{args.out}.pptx"
            prs.save(str(pptx_path))
            print(f"PPTX saved → {pptx_path}")

    browser.close()

print("Done.")
